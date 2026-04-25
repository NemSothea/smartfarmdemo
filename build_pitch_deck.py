"""SmartFarm investor pitch deck generator.

Produces SmartFarm-Pitch.pptx — 16 slides, 16:9, investor/judge tone,
English with Khmer terms (matches src/App.jsx).

Drop PNG/JPG screenshots into ./screenshots/ named:
  finance.png, calendar.png, guide.png, settings.png
and they'll be auto-embedded in the appendix slide.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand palette (mirrors AppColors in src/App.jsx) ─────────────
GREEN_700 = RGBColor(0x1F, 0x4D, 0x2E)
GREEN_600 = RGBColor(0x2E, 0x6E, 0x43)
GREEN_500 = RGBColor(0x4F, 0x9C, 0x6A)
GREEN_100 = RGBColor(0xD4, 0xE9, 0xDB)
GREEN_50  = RGBColor(0xE8, 0xF3, 0xEB)
SEC_600   = RGBColor(0x8B, 0x5A, 0x2B)
SEC_50    = RGBColor(0xF6, 0xEC, 0xDF)
DANGER    = RGBColor(0xB0, 0x38, 0x2A)
INK_900   = RGBColor(0x1A, 0x1F, 0x1C)
INK_700   = RGBColor(0x3A, 0x42, 0x3D)
INK_500   = RGBColor(0x6B, 0x6E, 0x6A)
INK_300   = RGBColor(0xB0, 0xB4, 0xB1)
INK_200   = RGBColor(0xD8, 0xDB, 0xD9)
INK_100   = RGBColor(0xEE, 0xF0, 0xEE)
SURFACE   = RGBColor(0xFF, 0xFF, 0xFF)
BG        = RGBColor(0xF5, 0xF6, 0xF5)

# ── Geometry ─────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────
def add_blank_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return slide


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK_700,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp


def add_pill(slide, x, y, w, h, label, fill, text_color, *, size=10, bold=False):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.shadow.inherit = False
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Emu(20000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return pill


def add_header(slide, eyebrow, title, *, accent=GREEN_600):
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.45), accent)
    add_text(slide, Inches(0.9), Inches(0.5), Inches(8), Inches(0.4),
             eyebrow.upper(), size=11, bold=True, color=GREEN_600)
    add_text(slide, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.7),
             title, size=28, bold=True, color=INK_900)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "SmartFarm — កសិកម្មឆ្លាតវៃ ងាយស្រួល",
             size=9, color=INK_500)
    add_text(slide, Inches(11.7), Inches(7.05), Inches(1.1), Inches(0.3),
             f"{page_num} / {total}", size=9, color=INK_500, align=PP_ALIGN.RIGHT)


# ── Phone-frame mockup (mirrors the 390×844 frame in src/App.jsx) ─
def add_phone_frame(slide, x, y, *, scale=0.0042):
    """Render a stylized phone frame at (x, y). Returns inner content rect."""
    pw = Emu(int(390 / 0.0042 * scale)) if False else Inches(2.8)
    ph = Inches(6.05)
    # Outer bezel
    bezel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pw, ph)
    bezel.adjustments[0] = 0.07
    bezel.shadow.inherit = False
    bezel.fill.solid()
    bezel.fill.fore_color.rgb = INK_900
    bezel.line.fill.background()
    # Inner screen
    inset = Inches(0.08)
    screen_x = x + inset
    screen_y = y + inset
    screen_w = pw - inset * 2
    screen_h = ph - inset * 2
    screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    screen_x, screen_y, screen_w, screen_h)
    screen.adjustments[0] = 0.05
    screen.shadow.inherit = False
    screen.fill.solid()
    screen.fill.fore_color.rgb = SURFACE
    screen.line.fill.background()
    # Notch
    notch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   screen_x + screen_w/2 - Inches(0.42),
                                   screen_y + Inches(0.06),
                                   Inches(0.84), Inches(0.18))
    notch.adjustments[0] = 0.5
    notch.shadow.inherit = False
    notch.fill.solid()
    notch.fill.fore_color.rgb = INK_900
    notch.line.fill.background()
    return (screen_x, screen_y, screen_w, screen_h)


def add_tab_bar(slide, sx, sy, sw, sh, tabs, active_idx):
    """Bottom tab bar matching src/App.jsx."""
    bar_h = Inches(0.55)
    bar_y = sy + sh - bar_h
    add_rect(slide, sx, bar_y - Pt(0.5), sw, Pt(0.5), INK_200)
    n = len(tabs)
    tab_w = sw / n
    for i, (icon, label) in enumerate(tabs):
        tx = sx + tab_w * i
        color = GREEN_600 if i == active_idx else INK_300
        add_text(slide, tx, bar_y + Inches(0.05), tab_w, Inches(0.25),
                 icon, size=12, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, tx, bar_y + Inches(0.3), tab_w, Inches(0.22),
                 label, size=7, bold=(i == active_idx),
                 color=color, align=PP_ALIGN.CENTER)


def add_app_chrome(slide, sx, sy, sw, sh, header_label, tabs, active_idx):
    """Status bar + page header + tab bar; returns the content rect (between header and tabs)."""
    # Status bar
    add_text(slide, sx + Inches(0.15), sy + Inches(0.08), Inches(0.5), Inches(0.2),
             "9:41", size=8, bold=True, color=INK_900)
    add_text(slide, sx + sw - Inches(0.5), sy + Inches(0.08), Inches(0.35), Inches(0.2),
             "▲▲▲", size=8, color=INK_900, align=PP_ALIGN.RIGHT)
    # Page header
    hy = sy + Inches(0.32)
    add_rect(slide, sx, hy, sw, Inches(0.4), SURFACE)
    add_rect(slide, sx, hy + Inches(0.4) - Pt(0.5), sw, Pt(0.5), INK_200)
    add_text(slide, sx + Inches(0.18), hy + Inches(0.08), sw, Inches(0.3),
             header_label, size=12, bold=True, color=INK_900)
    # Tab bar
    add_tab_bar(slide, sx, sy, sw, sh, tabs, active_idx)
    return (sx, hy + Inches(0.4),
            sw, sh - Inches(0.32) - Inches(0.4) - Inches(0.55))


# ── Content for each tab mockup ───────────────────────────────────
TABS = [
    ("$", "Finance"),
    ("=", "Calendar"),
    ("i", "Guide"),
    ("*", "Settings"),
]


def draw_finance_mock(slide, cx, cy, cw, ch):
    """Balance card + filter chips + 3 transaction rows."""
    pad = Inches(0.12)
    x = cx + pad
    w = cw - pad * 2
    # Balance card
    card_h = Inches(1.15)
    add_rect(slide, x, cy + pad, w, card_h, SURFACE, line=INK_200, radius=True)
    add_text(slide, x + Inches(0.15), cy + pad + Inches(0.08), w, Inches(0.18),
             "Balance / សមតុល្យ", size=8, color=INK_300)
    add_text(slide, x + Inches(0.15), cy + pad + Inches(0.24), w, Inches(0.4),
             "340,000 ៛", size=18, bold=True, color=GREEN_700)
    # KHR/USD pills
    add_pill(slide, x + w - Inches(0.95), cy + pad + Inches(0.18),
             Inches(0.4), Inches(0.22), "KHR", GREEN_50, GREEN_700, size=8, bold=True)
    add_pill(slide, x + w - Inches(0.5), cy + pad + Inches(0.18),
             Inches(0.4), Inches(0.22), "USD", SURFACE, INK_500, size=8)
    # Income / Expense split
    split_y = cy + pad + Inches(0.7)
    add_rect(slide, x + Inches(0.05), split_y, w/2 - Inches(0.05), Inches(0.4), SURFACE)
    add_text(slide, x + Inches(0.15), split_y, w/2, Inches(0.18),
             "● Income / ចំណូល", size=7, color=INK_500)
    add_text(slide, x + Inches(0.15), split_y + Inches(0.16), w/2, Inches(0.22),
             "430,000 ៛", size=10, bold=True, color=GREEN_600)
    add_text(slide, x + w/2 + Inches(0.05), split_y, w/2, Inches(0.18),
             "● Expense / ចំណាយ", size=7, color=INK_500)
    add_text(slide, x + w/2 + Inches(0.05), split_y + Inches(0.16), w/2, Inches(0.22),
             "90,000 ៛", size=10, bold=True, color=DANGER)
    # Filter chips
    chip_y = cy + pad + card_h + Inches(0.1)
    for i, (lbl, active) in enumerate([("All", True), ("Income", False), ("Expense", False)]):
        cx_pos = x + Inches(0.05) + Inches(0.65) * i
        fill = GREEN_100 if active else SURFACE
        textc = GREEN_700 if active else INK_500
        add_pill(slide, cx_pos, chip_y, Inches(0.6), Inches(0.25),
                 lbl, fill, textc, size=8, bold=active)
    # Transaction list
    list_y = chip_y + Inches(0.4)
    rows = [
        ("Sold rice / លក់ស្រូវ", "លក់ · 04-10", "+250,000 ៛", GREEN_600),
        ("Fertilizer / ទិញជី", "ជី · 04-08", "-45,000 ៛", DANGER),
        ("Labor / កម្មករ", "ការងារ · 04-06", "-30,000 ៛", DANGER),
    ]
    for i, (note, meta, amt, color) in enumerate(rows):
        ry = list_y + Inches(0.55) * i
        add_rect(slide, x, ry, w, Inches(0.5), SURFACE, line=INK_200, radius=True)
        add_text(slide, x + Inches(0.12), ry + Inches(0.07), w - Inches(1.3), Inches(0.2),
                 note, size=8, bold=True, color=INK_900)
        add_text(slide, x + Inches(0.12), ry + Inches(0.27), w, Inches(0.2),
                 meta, size=7, color=INK_300)
        add_text(slide, x + Inches(0.12), ry + Inches(0.13), w - Inches(0.25), Inches(0.25),
                 amt, size=9, bold=True, color=color, align=PP_ALIGN.RIGHT)
    # FAB
    fab = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx + cw - Inches(0.65), cy + ch - Inches(0.7),
                                 Inches(0.45), Inches(0.45))
    fab.shadow.inherit = False
    fab.fill.solid()
    fab.fill.fore_color.rgb = GREEN_600
    fab.line.fill.background()
    add_text(slide, cx + cw - Inches(0.65), cy + ch - Inches(0.78),
             Inches(0.45), Inches(0.45),
             "+", size=20, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def draw_calendar_mock(slide, cx, cy, cw, ch):
    pad = Inches(0.12)
    x = cx + pad
    w = cw - pad * 2
    # Calendar card
    cal_h = Inches(2.2)
    add_rect(slide, x, cy + pad, w, cal_h, SURFACE, line=INK_200, radius=True)
    add_text(slide, x, cy + pad + Inches(0.1), w, Inches(0.25),
             "មេសា 2026", size=11, bold=True, color=INK_900, align=PP_ALIGN.CENTER)
    # Day-of-week headers
    days = ["S", "M", "T", "W", "T", "F", "S"]
    grid_x = x + Inches(0.1)
    grid_w = w - Inches(0.2)
    cell_w = grid_w / 7
    head_y = cy + pad + Inches(0.4)
    for i, d in enumerate(days):
        add_text(slide, grid_x + cell_w * i, head_y, cell_w, Inches(0.18),
                 d, size=7, bold=True, color=INK_300, align=PP_ALIGN.CENTER)
    # Date grid (5 weeks for April 2026, starts Wed)
    grid_y = head_y + Inches(0.22)
    cell_h = Inches(0.28)
    activity_dates = {15, 12, 18, 25}
    selected = 15
    today = 25
    day = 1
    days_in_month = 30
    start_offset = 3  # April 1 is Wednesday
    for week in range(6):
        for col in range(7):
            idx = week * 7 + col
            if idx < start_offset or day > days_in_month:
                continue
            cell_x = grid_x + cell_w * col
            cell_y = grid_y + cell_h * week
            is_sel = day == selected
            is_today = day == today
            if is_sel:
                add_rect(slide, cell_x + Inches(0.04), cell_y, cell_w - Inches(0.08),
                         cell_h - Inches(0.04), GREEN_600, radius=True)
                color = SURFACE
            elif is_today:
                add_rect(slide, cell_x + Inches(0.04), cell_y, cell_w - Inches(0.08),
                         cell_h - Inches(0.04), GREEN_100, radius=True)
                color = GREEN_700
            else:
                color = INK_700
            add_text(slide, cell_x, cell_y + Inches(0.02), cell_w, Inches(0.2),
                     str(day), size=8, bold=is_sel, color=color, align=PP_ALIGN.CENTER)
            if day in activity_dates:
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                             cell_x + cell_w/2 - Inches(0.025),
                                             cell_y + Inches(0.21),
                                             Inches(0.05), Inches(0.05))
                dot.shadow.inherit = False
                dot.fill.solid()
                dot.fill.fore_color.rgb = SURFACE if is_sel else GREEN_500
                dot.line.fill.background()
            day += 1
    # Selected day list label
    list_y = cy + pad + cal_h + Inches(0.15)
    add_text(slide, x, list_y, w, Inches(0.22),
             "មេសា 15, 2026", size=10, bold=True, color=INK_700)
    # Activity rows
    acts = [
        ("Plant rice / ដាំស្រូវ", "ដំណាំ", False),
        ("Spray pesticide / បាញ់ថ្នាំ", "ការពារ", False),
    ]
    for i, (name, kind, done) in enumerate(acts):
        ry = list_y + Inches(0.3) + Inches(0.5) * i
        add_rect(slide, x, ry, w, Inches(0.45), SURFACE, line=INK_200, radius=True)
        add_text(slide, x + Inches(0.12), ry + Inches(0.06), w, Inches(0.22),
                 name, size=8, bold=True, color=INK_900)
        add_text(slide, x + Inches(0.12), ry + Inches(0.24), w, Inches(0.18),
                 kind, size=7, color=INK_300)
        check = "[x]" if done else "[ ]"
        add_text(slide, x + w - Inches(0.4), ry + Inches(0.12), Inches(0.3), Inches(0.25),
                 check, size=10, color=INK_500, align=PP_ALIGN.CENTER)


def draw_guide_mock(slide, cx, cy, cw, ch):
    pad = Inches(0.12)
    x = cx + pad
    w = cw - pad * 2
    # Search box
    add_rect(slide, x, cy + pad, w, Inches(0.32), SURFACE, line=INK_200, radius=True)
    add_text(slide, x + Inches(0.15), cy + pad + Inches(0.06), w, Inches(0.2),
             "ស្វែងរក...   (Search)", size=9, color=INK_300)
    # Category chips
    chip_y = cy + pad + Inches(0.45)
    chips = [("All", True), ("Insects / សត្វល្អិត", False),
             ("Fungal / ផ្សិត", False), ("Bacterial / បាក់តេរី", False)]
    cx_pos = x
    for lbl, active in chips:
        ww = Inches(0.55) if lbl == "All" else Inches(1.05)
        fill = GREEN_100 if active else SURFACE
        textc = GREEN_700 if active else INK_500
        add_pill(slide, cx_pos, chip_y, ww, Inches(0.25), lbl,
                 fill, textc, size=7, bold=active)
        cx_pos += ww + Inches(0.05)
    # Pest list
    list_y = chip_y + Inches(0.42)
    pests = [
        ("Rice worm / ដង្កូវស្រូវ", "Insects / សត្វល្អិត"),
        ("Leaf fungus / ផ្សិតស្លឹក", "Fungal / ផ្សិត"),
        ("Bacterial wilt / រោគប្លែក", "Bacterial / បាក់តេរី"),
        ("Rice planthopper / ចៃស្រូវ", "Insects / សត្វល្អិត"),
    ]
    for i, (name, kind) in enumerate(pests):
        ry = list_y + Inches(0.5) * i
        add_rect(slide, x, ry, w, Inches(0.45), SURFACE, line=INK_200, radius=True)
        add_text(slide, x + Inches(0.12), ry + Inches(0.06), w, Inches(0.22),
                 name, size=8, bold=True, color=INK_900)
        add_text(slide, x + Inches(0.12), ry + Inches(0.24), w, Inches(0.18),
                 kind, size=7, color=INK_300)
        add_text(slide, x + w - Inches(0.25), ry + Inches(0.1), Inches(0.2), Inches(0.25),
                 ">", size=12, color=INK_300)


def draw_settings_mock(slide, cx, cy, cw, ch):
    pad = Inches(0.12)
    x = cx + pad
    w = cw - pad * 2
    cards = [
        ("[T] Preferences / ចំណូលចិត្ត", [("Light / ភ្លឺ", True), ("Dark / ងងឹត", False)]),
        ("[#] Data Stats / ស្ថិតិទិន្នន័យ", [("5", "Transactions"), ("4", "Activities")]),
        ("[->] Export Data / នាំចេញ", [("CSV — Transactions", None)]),
        ("[B] Backup & Restore / បម្រុងទុក", [("Backup all (JSON)", None),
                                              ("Restore from file", None)]),
    ]
    cy_cur = cy + pad
    for title, items in cards:
        card_h = Inches(0.85) if title.startswith("[T]") else (
                 Inches(0.85) if title.startswith("[#]") else Inches(0.8))
        if title.startswith("[B]"):
            card_h = Inches(1.15)
        add_rect(slide, x, cy_cur, w, card_h, SURFACE, line=INK_200, radius=True)
        add_text(slide, x + Inches(0.12), cy_cur + Inches(0.08), w, Inches(0.25),
                 title, size=9, bold=True, color=INK_900)
        if title.startswith("[T]"):
            for i, (lbl, active) in enumerate(items):
                bx = x + Inches(0.1) + (w/2 - Inches(0.05)) * i
                bw = w/2 - Inches(0.15)
                fill = GREEN_50 if active else SURFACE
                textc = GREEN_700 if active else INK_500
                add_rect(slide, bx, cy_cur + Inches(0.4), bw, Inches(0.3),
                         fill, line=GREEN_600 if active else INK_200, radius=True)
                add_text(slide, bx, cy_cur + Inches(0.43), bw, Inches(0.25),
                         lbl, size=8, bold=active, color=textc, align=PP_ALIGN.CENTER)
        elif title.startswith("[#]"):
            for i, (n, lbl) in enumerate(items):
                bx = x + Inches(0.1) + (w/2 - Inches(0.05)) * i
                bw = w/2 - Inches(0.15)
                add_rect(slide, bx, cy_cur + Inches(0.38), bw, Inches(0.4),
                         BG, radius=True)
                add_text(slide, bx, cy_cur + Inches(0.4), bw, Inches(0.22),
                         n, size=12, bold=True, color=GREEN_700, align=PP_ALIGN.CENTER)
                add_text(slide, bx, cy_cur + Inches(0.6), bw, Inches(0.18),
                         lbl, size=7, color=INK_500, align=PP_ALIGN.CENTER)
        else:
            for i, (lbl, _) in enumerate(items):
                ry = cy_cur + Inches(0.36) + Inches(0.32) * i
                add_rect(slide, x + Inches(0.1), ry, w - Inches(0.2),
                         Inches(0.28), SURFACE, line=INK_200, radius=True)
                add_text(slide, x + Inches(0.2), ry + Inches(0.05),
                         w, Inches(0.22),
                         lbl, size=8, color=INK_700)
        cy_cur += card_h + Inches(0.1)


def draw_phone_with(slide, x, y, header_label, active_idx, draw_content):
    sx, sy, sw, sh = add_phone_frame(slide, x, y)
    cx, cy, cw, ch = add_app_chrome(slide, sx, sy, sw, sh, header_label, TABS, active_idx)
    draw_content(slide, cx, cy, cw, ch)


# ── Slide builders ───────────────────────────────────────────────
TOTAL_SLIDES = 16


def slide_title(prs):
    slide = add_blank_slide(prs)
    # Hero band
    add_rect(slide, 0, 0, SLIDE_W, Inches(7.5), SURFACE)
    add_rect(slide, 0, 0, Inches(5.5), SLIDE_H, GREEN_600)
    # Logo
    add_text(slide, Inches(0.6), Inches(2.5), Inches(4.3), Inches(1.6),
             "🌱", size=110, color=SURFACE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.2), Inches(4.3), Inches(0.6),
             "SmartFarm", size=36, bold=True,
             color=SURFACE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.8), Inches(4.3), Inches(0.4),
             "កសិកម្មឆ្លាតវៃ ងាយស្រួល",
             size=14, color=SURFACE, align=PP_ALIGN.CENTER)
    # Right side
    add_text(slide, Inches(6.2), Inches(2.0), Inches(6.5), Inches(0.4),
             "INVESTOR PITCH · 2026", size=11, bold=True, color=GREEN_600)
    add_text(slide, Inches(6.2), Inches(2.5), Inches(6.5), Inches(2.2),
             "An offline farm-management app\nbuilt for Cambodian small-scale farmers.",
             size=26, bold=True, color=INK_900)
    add_text(slide, Inches(6.2), Inches(4.4), Inches(6.5), Inches(1.5),
             "Track money. Plan activities.\nIdentify pests. All in Khmer. All offline.",
             size=14, color=INK_500)
    add_text(slide, Inches(6.2), Inches(6.6), Inches(6.5), Inches(0.4),
             "Presented by the SmartFarm team",
             size=10, color=INK_300)


def slide_problem(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "The Problem", "Farming small plots is a daily juggle of memory and guesswork")
    items = [
        ("[$]", "Untracked money",
         "Income and expenses live in notebooks or memory.\n"
         "By harvest time, no one knows real profit."),
        ("[=]", "Missed activities",
         "Planting, watering, spraying, harvesting — easy to lose the schedule.\n"
         "A missed window costs a season."),
        ("[!]", "Unknown pests & disease",
         "By the time damage is visible, the right treatment\n"
         "is often guessed at or asked from neighbors."),
        ("[X]", "Patchy connectivity",
         "Rural areas have weak or no internet.\n"
         "Cloud-only farming apps don't work where farmers actually farm."),
    ]
    cols = 2
    cw = Inches(5.7)
    ch = Inches(2.4)
    gap = Inches(0.4)
    sx = Inches(0.7)
    sy = Inches(1.7)
    for i, (icon, title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = sx + (cw + gap) * col
        y = sy + (ch + Inches(0.15)) * row
        add_rect(slide, x, y, cw, ch, SURFACE, line=INK_200, radius=True)
        add_rect(slide, x, y, Inches(0.18), ch, GREEN_500)
        add_text(slide, x + Inches(0.4), y + Inches(0.2), Inches(1), Inches(0.4),
                 icon, size=20, bold=True, color=GREEN_600)
        add_text(slide, x + Inches(1.3), y + Inches(0.25), cw - Inches(1.3), Inches(0.5),
                 title, size=16, bold=True, color=INK_900)
        add_text(slide, x + Inches(1.3), y + Inches(0.85), cw - Inches(1.5), ch - Inches(1),
                 body, size=11, color=INK_500)
    add_footer(slide, 2, TOTAL_SLIDES)


def slide_solution(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "The Solution", "One offline app, four tools farmers already need")
    blurbs = [
        ("[$]", "Finance Tracker",
         "Income & expense in KHR or USD.\nLive balance.\nCategories, filters, dual currency."),
        ("[=]", "Calendar & Activities",
         "Month grid with activity dots.\nPlant / Tend / Protect / Harvest types.\nMark done, never forget."),
        ("[i]", "Pest & Disease Guide",
         "Bundled offline library.\nSearch + category filters.\nSymptoms, treatment, prevention."),
        ("[*]", "Settings & Backup",
         "Light / dark theme.\nCSV export, JSON backup & restore.\nOne tap to clear & start fresh."),
    ]
    sx = Inches(0.7)
    sy = Inches(1.7)
    cw = Inches(2.95)
    ch = Inches(4.0)
    gap = Inches(0.15)
    for i, (icon, title, body) in enumerate(blurbs):
        x = sx + (cw + gap) * i
        add_rect(slide, x, sy, cw, ch, SURFACE, line=INK_200, radius=True)
        add_rect(slide, x, sy, cw, Inches(1.0), GREEN_50)
        add_text(slide, x, sy + Inches(0.2), cw, Inches(0.6),
                 icon, size=28, bold=True, color=GREEN_600, align=PP_ALIGN.CENTER)
        add_text(slide, x, sy + Inches(1.15), cw, Inches(0.4),
                 title, size=14, bold=True, color=INK_900, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.25), sy + Inches(1.65), cw - Inches(0.5), Inches(2.2),
                 body, size=10, color=INK_500, align=PP_ALIGN.CENTER)
    # Bottom strip
    add_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.7),
             GREEN_50, radius=True)
    add_text(slide, Inches(0.9), Inches(6.18), Inches(11.5), Inches(0.4),
             "100% offline   ·   Khmer-first UI   ·   KHR + USD   ·   No accounts, no data leaves the phone",
             size=12, bold=True, color=GREEN_700, align=PP_ALIGN.CENTER)
    add_footer(slide, 3, TOTAL_SLIDES)


def slide_user_flow(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "User Flow", "From first launch to daily use")
    # Flow boxes
    steps = [
        ("Splash", "🌱  Brand intro\n2 seconds"),
        ("Onboarding", "3 pages\nFinance · Plan · Reports"),
        ("Main App", "Tab bar with\n4 modules"),
        ("Daily Use", "Add txn · Plan task\nLook up pest · Backup"),
    ]
    sx = Inches(0.7)
    sy = Inches(2.4)
    cw = Inches(2.6)
    ch = Inches(2.0)
    gap_x = (Inches(11.9) - cw * 4) / 3
    for i, (title, body) in enumerate(steps):
        x = sx + (cw + gap_x) * i
        add_rect(slide, x, sy, cw, ch, SURFACE, line=INK_200, radius=True)
        add_text(slide, x, sy + Inches(0.25), cw, Inches(0.5),
                 title, size=18, bold=True, color=GREEN_700, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.25), sy + Inches(0.95), cw - Inches(0.5),
                 ch - Inches(1), body, size=12, color=INK_500, align=PP_ALIGN.CENTER)
        # Arrow between steps
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           x + cw + Inches(0.05),
                                           sy + ch/2 - Inches(0.18),
                                           gap_x - Inches(0.1),
                                           Inches(0.36))
            arrow.shadow.inherit = False
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREEN_500
            arrow.line.fill.background()
    # Caption
    add_text(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(0.5),
             "First run takes ~30 seconds.   Daily use is single-tap actions.",
             size=14, bold=True, color=INK_900, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.4),
             "Persistent local storage means the app opens with your data — every time, online or off.",
             size=11, color=INK_500, align=PP_ALIGN.CENTER)
    add_footer(slide, 4, TOTAL_SLIDES)


def slide_finance(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Feature 1 / 4", "Finance Tracker — track every riel and dollar")
    # Phone mockup on the left
    draw_phone_with(slide, Inches(0.7), Inches(1.3),
                    "[$] ហិរញ្ញវត្ថុ (Finance)", 0, draw_finance_mock)
    # Right content
    rx = Inches(4.2)
    ry = Inches(1.7)
    add_text(slide, rx, ry, Inches(8), Inches(0.4),
             "Why farmers care", size=13, bold=True, color=GREEN_600)
    add_text(slide, rx, ry + Inches(0.45), Inches(8.5), Inches(1.2),
             "Profit isn't profit until you've counted it.\n"
             "Most small farmers know revenue but not net.",
             size=14, color=INK_700)
    add_text(slide, rx, ry + Inches(1.85), Inches(8), Inches(0.4),
             "What it does", size=13, bold=True, color=GREEN_600)
    bullets = [
        "•  Add income / expense in seconds (FAB → bottom sheet)",
        "•  Dual currency: KHR ↔ USD at a stable rate (4,100)",
        "•  Categories: ជី, ពូជ, ការងារ, ឧបករណ៍, លក់",
        "•  Live balance, income, expense breakdown",
        "•  Filter chips: All / Income / Expense",
        "•  CSV export for record-keeping or microfinance applications",
    ]
    add_text(slide, rx, ry + Inches(2.3), Inches(8.5), Inches(3),
             "\n".join(bullets), size=12, color=INK_700)
    # Process flow ribbon
    add_rect(slide, Inches(4.2), Inches(6.0), Inches(8.5), Inches(0.8),
             GREEN_50, radius=True)
    add_text(slide, Inches(4.4), Inches(6.05), Inches(8), Inches(0.3),
             "Add Transaction Flow", size=10, bold=True, color=GREEN_700)
    add_text(slide, Inches(4.4), Inches(6.32), Inches(8), Inches(0.45),
             "Tap +  →  Income / Expense  →  Amount + Date + Currency  →  Category + Note  →  Save",
             size=11, color=INK_700)
    add_footer(slide, 5, TOTAL_SLIDES)


def slide_calendar(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Feature 2 / 4", "Calendar & Activities — never miss a season")
    draw_phone_with(slide, Inches(0.7), Inches(1.3),
                    "[=] ប្រតិទិន (Calendar)", 1, draw_calendar_mock)
    rx = Inches(4.2)
    ry = Inches(1.7)
    add_text(slide, rx, ry, Inches(8), Inches(0.4),
             "Why farmers care", size=13, bold=True, color=GREEN_600)
    add_text(slide, rx, ry + Inches(0.45), Inches(8.5), Inches(1.2),
             "Crop cycles are unforgiving.\n"
             "A missed spray window or a late harvest costs an entire season's margin.",
             size=14, color=INK_700)
    add_text(slide, rx, ry + Inches(1.85), Inches(8), Inches(0.4),
             "What it does", size=13, bold=True, color=GREEN_600)
    bullets = [
        "•  Month-grid view with dot markers on busy days",
        "•  Activity types: ដំណាំ (Plant) · ថែទាំ (Tend) · ការពារ (Protect) · ប្រមូល (Harvest)",
        "•  Tap a day → see only that day's activities",
        "•  Mark done with a single tap",
        "•  Add → Name + Date + Type → save → calendar dot appears",
    ]
    add_text(slide, rx, ry + Inches(2.3), Inches(8.5), Inches(2.5),
             "\n".join(bullets), size=12, color=INK_700)
    # Flow ribbon
    add_rect(slide, Inches(4.2), Inches(6.0), Inches(8.5), Inches(0.8),
             GREEN_50, radius=True)
    add_text(slide, Inches(4.4), Inches(6.05), Inches(8), Inches(0.3),
             "Schedule Activity Flow", size=10, bold=True, color=GREEN_700)
    add_text(slide, Inches(4.4), Inches(6.32), Inches(8), Inches(0.45),
             "Pick day  →  Tap +  →  Name + Date + Type  →  Save  →  Dot on calendar",
             size=11, color=INK_700)
    add_footer(slide, 6, TOTAL_SLIDES)


def slide_guide(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Feature 3 / 4", "Pest & Disease Guide — answers without internet")
    draw_phone_with(slide, Inches(0.7), Inches(1.3),
                    "[i] មគ្គុទ្ទេសក៍ (Guide)", 2, draw_guide_mock)
    rx = Inches(4.2)
    ry = Inches(1.7)
    add_text(slide, rx, ry, Inches(8), Inches(0.4),
             "Why farmers care", size=13, bold=True, color=GREEN_600)
    add_text(slide, rx, ry + Inches(0.45), Inches(8.5), Inches(1.2),
             "When a pest appears, decisions are minutes — not days.\n"
             "Online services aren't an option in the field.",
             size=14, color=INK_700)
    add_text(slide, rx, ry + Inches(1.85), Inches(8), Inches(0.4),
             "What it does", size=13, bold=True, color=GREEN_600)
    bullets = [
        "•  Bundled offline library — no network call ever",
        "•  Search by Khmer name",
        "•  Category filter: សត្វល្អិត (Insects), ផ្សិត (Fungal), បាក់តេរី (Bacterial)",
        "•  Detail view with expandable sections:",
        "      រោគសញ្ញា (Symptoms) · វិធីព្យាបាល (Treatment) · វិធីការពារ (Prevention)",
        "•  Designed to scale to 30–50 entries with photos",
    ]
    add_text(slide, rx, ry + Inches(2.3), Inches(8.5), Inches(2.5),
             "\n".join(bullets), size=12, color=INK_700)
    # Flow ribbon
    add_rect(slide, Inches(4.2), Inches(6.0), Inches(8.5), Inches(0.8),
             GREEN_50, radius=True)
    add_text(slide, Inches(4.4), Inches(6.05), Inches(8), Inches(0.3),
             "Lookup Flow", size=10, bold=True, color=GREEN_700)
    add_text(slide, Inches(4.4), Inches(6.32), Inches(8), Inches(0.45),
             "Search or Filter  →  Tap pest  →  Expand sections  →  Apply treatment",
             size=11, color=INK_700)
    add_footer(slide, 7, TOTAL_SLIDES)


def slide_settings(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Feature 4 / 4", "Settings & Backup — own your data")
    draw_phone_with(slide, Inches(0.7), Inches(1.3),
                    "[*] ការកំណត់ (Settings)", 3, draw_settings_mock)
    rx = Inches(4.2)
    ry = Inches(1.7)
    add_text(slide, rx, ry, Inches(8), Inches(0.4),
             "Why farmers care", size=13, bold=True, color=GREEN_600)
    add_text(slide, rx, ry + Inches(0.45), Inches(8.5), Inches(1.2),
             "Phones break. Phones get sold. Phones get reset.\n"
             "Two seasons of records vanish in one moment without backup.",
             size=14, color=INK_700)
    add_text(slide, rx, ry + Inches(1.85), Inches(8), Inches(0.4),
             "What it does", size=13, bold=True, color=GREEN_600)
    bullets = [
        "•  Light / dark theme — usable in bright sun and at night",
        "•  Data stats — see counts at a glance",
        "•  CSV export — share with cooperatives, MFIs, tax officers",
        "•  Full JSON backup — one tap, save to Files / share / cloud",
        "•  Restore — load a backup file, replace current data",
        "•  Clear all — fresh start, with a confirmation guard",
    ]
    add_text(slide, rx, ry + Inches(2.3), Inches(8.5), Inches(2.5),
             "\n".join(bullets), size=12, color=INK_700)
    # Flow ribbon
    add_rect(slide, Inches(4.2), Inches(6.0), Inches(8.5), Inches(0.8),
             GREEN_50, radius=True)
    add_text(slide, Inches(4.4), Inches(6.05), Inches(8), Inches(0.3),
             "Backup & Restore Flow", size=10, bold=True, color=GREEN_700)
    add_text(slide, Inches(4.4), Inches(6.32), Inches(8), Inches(0.45),
             "Settings  →  Backup all (JSON)  →  Save file  →  (later) Restore from file  →  Confirm",
             size=11, color=INK_700)
    add_footer(slide, 8, TOTAL_SLIDES)


def slide_differentiators(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Why Now & Why Us", "Four reasons SmartFarm is different")
    items = [
        ("[OFF]", "Truly offline",
         "Every feature works with airplane mode on.\n"
         "Reference data is bundled. Storage is local.\n"
         "Cloud-first competitors fail in rural Cambodia."),
        ("[KH]", "Khmer-first by design",
         "Every label, every category, every pest name —\n"
         "in Khmer. Not a translation layer."),
        ("[$$]", "Dual currency native",
         "KHR and USD both first-class.\n"
         "Reflects how Cambodian commerce actually works."),
        ("[--]", "No account, no friction",
         "Open the app. Use it.\n"
         "No login, no email, no data leaves the phone."),
    ]
    sx = Inches(0.7)
    sy = Inches(1.7)
    cw = Inches(5.95)
    ch = Inches(2.5)
    for i, (icon, title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = sx + (cw + Inches(0.2)) * col
        y = sy + (ch + Inches(0.15)) * row
        add_rect(slide, x, y, cw, ch, SURFACE, line=INK_200, radius=True)
        add_rect(slide, x, y, Inches(0.18), ch, GREEN_600)
        add_text(slide, x + Inches(0.4), y + Inches(0.2), Inches(1.6), Inches(0.5),
                 icon, size=18, bold=True, color=GREEN_600)
        add_text(slide, x + Inches(2.0), y + Inches(0.25), cw - Inches(2.2), Inches(0.5),
                 title, size=16, bold=True, color=INK_900)
        add_text(slide, x + Inches(2.0), y + Inches(0.85), cw - Inches(2.3), ch - Inches(1),
                 body, size=11, color=INK_500)
    add_footer(slide, 9, TOTAL_SLIDES)


def slide_status(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Status & Tech Stack", "Working demo today, native iOS tomorrow")
    # Left: status cards
    sx = Inches(0.7)
    sy = Inches(1.7)
    cards = [
        ("Today", GREEN_600,
         "•  Working React + Vite demo\n"
         "•  Khmer UI, KHR/USD, theme toggle\n"
         "•  All four modules functional\n"
         "•  CSV export + JSON backup live"),
        ("Next 3 months", SEC_600,
         "•  Native iOS port (Swift / SwiftUI)\n"
         "•  iOS 13+ deployment target\n"
         "•  CoreData persistence\n"
         "•  Local notifications for activities"),
        ("6–12 months", GREEN_700,
         "•  Photo-based pest detail (30–50 entries)\n"
         "•  Voice notes for low-literacy users\n"
         "•  TestFlight pilot with real farmers\n"
         "•  Optional iCloud Drive backup"),
    ]
    cw = Inches(3.95)
    ch = Inches(4.6)
    for i, (label, color, body) in enumerate(cards):
        x = sx + (cw + Inches(0.15)) * i
        add_rect(slide, x, sy, cw, ch, SURFACE, line=INK_200, radius=True)
        add_rect(slide, x, sy, cw, Inches(0.55), color)
        add_text(slide, x, sy + Inches(0.12), cw, Inches(0.4),
                 label, size=14, bold=True, color=SURFACE, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), sy + Inches(0.85), cw - Inches(0.5), ch - Inches(1),
                 body, size=11, color=INK_700)
    # Bottom: stack pill row
    add_rect(slide, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5),
             INK_100, radius=True)
    add_text(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4),
             "Stack:  React 18  ·  Vite  ·  localStorage   →   SwiftUI  ·  CoreData  ·  iOS 13+",
             size=11, bold=True, color=INK_700, align=PP_ALIGN.CENTER)
    add_footer(slide, 10, TOTAL_SLIDES)


def slide_competitive(prs):
    """Competitive landscape — comparison table.

    Honest framing: most Cambodian farmers' real alternative isn't another app.
    It's a notebook. The other "real" option is a generic cloud app.
    """
    slide = add_blank_slide(prs)
    add_header(slide, "Competitive Landscape",
               "Why farmers will pick SmartFarm over what they have today")
    # Column headers
    table_x = Inches(0.7)
    table_y = Inches(1.7)
    feature_w = Inches(4.0)
    col_w = Inches(2.55)
    row_h = Inches(0.45)
    cols = [
        ("Paper notebook", INK_500, "(today)"),
        ("Cloud farming app", SEC_600, "e.g. Plantix"),
        ("SmartFarm", GREEN_700, "(us)"),
    ]
    # Header row
    add_text(slide, table_x, table_y, feature_w, row_h,
             "Feature", size=11, bold=True, color=INK_700,
             anchor=MSO_ANCHOR.MIDDLE)
    for i, (name, color, sub) in enumerate(cols):
        cx = table_x + feature_w + col_w * i
        is_us = (name == "SmartFarm")
        if is_us:
            add_rect(slide, cx, table_y, col_w, row_h, GREEN_50, radius=True)
        add_text(slide, cx, table_y + Inches(0.04), col_w, Inches(0.22),
                 name, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, cx, table_y + Inches(0.24), col_w, Inches(0.18),
                 sub, size=8, color=INK_300, align=PP_ALIGN.CENTER)
    # Feature rows: True = green check, False = red X, "~" = partial
    rows = [
        ("Works fully offline",            True,  False, True),
        ("Khmer-first interface",          True,  False, True),
        ("Searchable records",             False, True,  True),
        ("Built-in pest & disease library", False, True,  True),
        ("Activity reminders & calendar",  False, True,  True),
        ("Backup & restore",               False, True,  True),
        ("Free, no account required",      True,  False, True),
        ("Designed for Cambodian farmers", "~",   False, True),
    ]
    for r, (label, *vals) in enumerate(rows):
        ry = table_y + row_h + Inches(0.05) + (row_h * r)
        # Zebra
        if r % 2 == 0:
            add_rect(slide, table_x, ry, feature_w + col_w * 3,
                     row_h, INK_100)
        add_text(slide, table_x + Inches(0.15), ry, feature_w, row_h,
                 label, size=10, color=INK_700, anchor=MSO_ANCHOR.MIDDLE)
        for i, v in enumerate(vals):
            cx = table_x + feature_w + col_w * i
            if v is True:
                mark, color = "✓", GREEN_600
            elif v is False:
                mark, color = "✗", DANGER
            else:
                mark, color = "~", SEC_600
            add_text(slide, cx, ry, col_w, row_h, mark,
                     size=18, bold=True, color=color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bottom callout
    add_rect(slide, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.55),
             GREEN_600, radius=True)
    add_text(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5),
             "SmartFarm is the only option that wins on every axis "
             "rural Cambodian farmers actually care about.",
             size=12, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 12, TOTAL_SLIDES)


def slide_team(prs):
    """Team — placeholder cards. EDIT NAMES/ROLES BEFORE PITCHING."""
    slide = add_blank_slide(prs)
    add_header(slide, "The Team", "Builders close to the problem")
    members = [
        ("ST", "[Founder Name]", "Founder · Lead Engineer",
         "Builds the app end-to-end.\n"
         "Background in iOS / mobile engineering.\n"
         "Connection to Cambodian farming community."),
        ("?", "[Co-founder Name]", "Product & Design",
         "Owns the Khmer UX and farmer research.\n"
         "Translates fieldwork into product.\n"
         "Replace with real bio."),
        ("?", "[Domain Advisor]", "Agriculture Advisor",
         "Real expertise in Cambodian small-scale farming.\n"
         "Validates pest data, activity types, terminology.\n"
         "Replace with real bio."),
    ]
    sx = Inches(0.7)
    sy = Inches(1.8)
    cw = Inches(3.95)
    ch = Inches(4.4)
    gap = Inches(0.15)
    for i, (initials, name, role, bio) in enumerate(members):
        x = sx + (cw + gap) * i
        # Card
        add_rect(slide, x, sy, cw, ch, SURFACE, line=INK_200, radius=True)
        # Avatar
        avatar_d = Inches(1.5)
        ax = x + (cw - avatar_d) / 2
        ay = sy + Inches(0.4)
        avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, ax, ay,
                                        avatar_d, avatar_d)
        avatar.shadow.inherit = False
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = GREEN_600 if i == 0 else (
            SEC_600 if i == 1 else GREEN_500)
        avatar.line.fill.background()
        add_text(slide, ax, ay, avatar_d, avatar_d, initials,
                 size=44, bold=True, color=SURFACE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name
        add_text(slide, x + Inches(0.2), sy + Inches(2.15),
                 cw - Inches(0.4), Inches(0.4),
                 name, size=16, bold=True, color=INK_900,
                 align=PP_ALIGN.CENTER)
        # Role
        add_text(slide, x + Inches(0.2), sy + Inches(2.55),
                 cw - Inches(0.4), Inches(0.3),
                 role, size=11, bold=True, color=GREEN_600,
                 align=PP_ALIGN.CENTER)
        # Divider
        add_rect(slide, x + Inches(1.3), sy + Inches(2.95),
                 cw - Inches(2.6), Inches(0.03), INK_200)
        # Bio
        add_text(slide, x + Inches(0.3), sy + Inches(3.1),
                 cw - Inches(0.6), Inches(1.2),
                 bio, size=11, color=INK_500, align=PP_ALIGN.CENTER)
    # Edit-warning footer note
    add_text(slide, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.4),
             "Note: names and bios above are placeholders — edit before pitching.",
             size=10, color=INK_300, align=PP_ALIGN.CENTER)
    add_footer(slide, 13, TOTAL_SLIDES)


def slide_ask(prs):
    """The Ask — funding amount, use of funds, milestones.

    Numbers below are PLACEHOLDERS — edit before pitching.
    """
    slide = add_blank_slide(prs)
    add_header(slide, "The Ask", "What we're raising and what it builds")
    # Big funding figure (left)
    add_rect(slide, Inches(0.7), Inches(1.7), Inches(4.6), Inches(4.6),
             GREEN_600, radius=True)
    add_text(slide, Inches(0.7), Inches(2.0), Inches(4.6), Inches(0.5),
             "SEED ROUND", size=12, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.7), Inches(2.6), Inches(4.6), Inches(1.2),
             "$50,000", size=64, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.7), Inches(3.9), Inches(4.6), Inches(0.5),
             "USD · 12 months runway",
             size=14, color=SURFACE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.2), Inches(4.7), Inches(1.6), Inches(0.04), SURFACE)
    add_text(slide, Inches(0.9), Inches(4.85), Inches(4.2), Inches(1.4),
             "From idea to pilot:\n"
             "real farmers, real iOS app,\nreal data — in one year.",
             size=13, color=SURFACE, align=PP_ALIGN.CENTER)
    # Right column: use of funds
    rx = Inches(5.6)
    ry = Inches(1.7)
    add_text(slide, rx, ry, Inches(7.2), Inches(0.4),
             "Use of funds", size=14, bold=True, color=GREEN_600)
    breakdown = [
        ("Engineering — iOS native port", 40, "$20,000"),
        ("Content — pest library (50 entries + photos)", 20, "$10,000"),
        ("Pilot — TestFlight cohort + farmer support", 20, "$10,000"),
        ("Design & localization polish", 12, "$6,000"),
        ("Operations & contingency", 8, "$4,000"),
    ]
    bar_x0 = rx
    bar_w_max = Inches(4.4)
    by = ry + Inches(0.5)
    for i, (label, pct, amt) in enumerate(breakdown):
        row_y = by + Inches(0.6) * i
        add_text(slide, bar_x0, row_y, Inches(5.0), Inches(0.22),
                 label, size=10, bold=True, color=INK_900)
        add_text(slide, bar_x0 + Inches(5.0), row_y, Inches(1.2), Inches(0.22),
                 amt, size=10, bold=True, color=GREEN_700, align=PP_ALIGN.RIGHT)
        add_text(slide, bar_x0 + Inches(6.3), row_y, Inches(0.7), Inches(0.22),
                 f"{pct}%", size=10, color=INK_500, align=PP_ALIGN.RIGHT)
        # Bar background
        add_rect(slide, bar_x0, row_y + Inches(0.28), bar_w_max,
                 Inches(0.16), INK_100, radius=True)
        # Bar fill
        add_rect(slide, bar_x0, row_y + Inches(0.28),
                 Emu(int(bar_w_max * pct / 100)),
                 Inches(0.16), GREEN_500, radius=True)
    # Milestones strip
    add_rect(slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.6),
             GREEN_50, radius=True)
    add_text(slide, Inches(0.9), Inches(6.45), Inches(2.3), Inches(0.5),
             "MILESTONES", size=10, bold=True, color=GREEN_700,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(2.5), Inches(6.45), Inches(10), Inches(0.5),
             "M3: iOS alpha   →   M6: 50-pest library shipped   →   "
             "M9: TestFlight pilot (20 farmers)   →   M12: public release",
             size=11, color=INK_700, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 14, TOTAL_SLIDES)


def slide_traction(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "The Opportunity", "A real audience, a real product, a clear path")
    # Three big metric tiles
    sx = Inches(0.7)
    sy = Inches(1.8)
    cw = Inches(3.95)
    ch = Inches(2.4)
    tiles = [
        ("Cambodia", "agriculture-led economy",
         "Small farmers form the backbone of rural livelihoods.\n"
         "Few have any digital tool today."),
        ("Smartphones", "wide & growing reach",
         "Entry-level Android and older iPhones are common,\n"
         "even in rural districts."),
        ("Connectivity", "still the blocker",
         "Network coverage is uneven; data is costly.\n"
         "Offline-first is the unlock."),
    ]
    for i, (h1, h2, body) in enumerate(tiles):
        x = sx + (cw + Inches(0.15)) * i
        add_rect(slide, x, sy, cw, ch, SURFACE, line=INK_200, radius=True)
        add_text(slide, x + Inches(0.3), sy + Inches(0.2), cw, Inches(0.5),
                 h1, size=22, bold=True, color=GREEN_700)
        add_text(slide, x + Inches(0.3), sy + Inches(0.75), cw, Inches(0.4),
                 h2, size=12, bold=True, color=SEC_600)
        add_text(slide, x + Inches(0.3), sy + Inches(1.2), cw - Inches(0.5), ch - Inches(1.3),
                 body, size=11, color=INK_500)
    # Bottom callouts
    add_rect(slide, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.0),
             GREEN_50, radius=True)
    add_text(slide, Inches(0.95), Inches(4.8), Inches(11.5), Inches(0.4),
             "What we have right now", size=12, bold=True, color=GREEN_700)
    add_text(slide, Inches(0.95), Inches(5.2), Inches(11.5), Inches(1.4),
             "•  A working app, not a deck — every feature on the previous slides is built and runs offline today.\n"
             "•  Khmer UX validated against the real domain (categories, activity types, common pests).\n"
             "•  A clear, conservative roadmap to a native iOS app supporting devices farmers actually own.",
             size=12, color=INK_700)
    add_footer(slide, 11, TOTAL_SLIDES)


def slide_appendix_screenshots(prs):
    """Appendix slide — 4 phone-shaped screenshot slots.

    Auto-embeds PNG/JPG from ./screenshots/ if files exist:
      finance.png, calendar.png, guide.png, settings.png
    Otherwise, draws a placeholder with the expected filename.
    """
    slide = add_blank_slide(prs)
    add_header(slide, "Appendix", "Live UI screenshots")
    slots = [
        ("Finance",  "finance",  "[$] ហិរញ្ញវត្ថុ"),
        ("Calendar", "calendar", "[=] ប្រតិទិន"),
        ("Guide",    "guide",    "[i] មគ្គុទ្ទេសក៍"),
        ("Settings", "settings", "[*] ការកំណត់"),
    ]
    sy = Inches(1.7)
    pw = Inches(2.55)
    ph = Inches(4.9)
    total_w = pw * 4 + Inches(0.3) * 3
    sx = (SLIDE_W - total_w) / 2
    screenshots_dir = os.path.join(os.path.dirname(__file__), "screenshots")
    for i, (label, key, kh) in enumerate(slots):
        x = sx + (pw + Inches(0.3)) * i
        # Caption above
        add_text(slide, x, sy - Inches(0.45), pw, Inches(0.25),
                 label, size=12, bold=True, color=INK_900,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x, sy - Inches(0.22), pw, Inches(0.22),
                 kh, size=9, color=INK_500, align=PP_ALIGN.CENTER)
        # Phone bezel
        bezel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sy, pw, ph)
        bezel.adjustments[0] = 0.06
        bezel.shadow.inherit = False
        bezel.fill.solid()
        bezel.fill.fore_color.rgb = INK_900
        bezel.line.fill.background()
        # Inner screen
        inset = Inches(0.07)
        inner_x = x + inset
        inner_y = sy + inset
        inner_w = pw - inset * 2
        inner_h = ph - inset * 2
        # Try each common extension
        embedded = None
        for ext in ("png", "jpg", "jpeg"):
            candidate = os.path.join(screenshots_dir, f"{key}.{ext}")
            if os.path.isfile(candidate):
                embedded = candidate
                break
        if embedded:
            slide.shapes.add_picture(embedded, inner_x, inner_y,
                                     width=inner_w, height=inner_h)
        else:
            # Placeholder
            placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                 inner_x, inner_y,
                                                 inner_w, inner_h)
            placeholder.adjustments[0] = 0.04
            placeholder.shadow.inherit = False
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = INK_100
            placeholder.line.color.rgb = INK_200
            placeholder.line.width = Pt(0.75)
            add_text(slide, inner_x, inner_y + inner_h/2 - Inches(0.6),
                     inner_w, Inches(0.4),
                     "[+]", size=24, color=INK_300,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, inner_x, inner_y + inner_h/2 - Inches(0.1),
                     inner_w, Inches(0.4),
                     "Drop screenshot at",
                     size=10, color=INK_500, align=PP_ALIGN.CENTER)
            add_text(slide, inner_x, inner_y + inner_h/2 + Inches(0.25),
                     inner_w, Inches(0.4),
                     f"screenshots/{key}.png",
                     size=10, bold=True, color=GREEN_700,
                     align=PP_ALIGN.CENTER, font="Consolas")
    # Footer hint
    add_text(slide, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.3),
             "Tip: rerun  python3 build_pitch_deck.py  after adding screenshots — "
             "they'll embed automatically.",
             size=10, color=INK_500, align=PP_ALIGN.CENTER, font="Consolas")
    add_footer(slide, 15, TOTAL_SLIDES)


def slide_close(prs):
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GREEN_600)
    add_text(slide, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.0),
             "Thank you", size=64, bold=True, color=SURFACE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(3.2), Inches(12.1), Inches(0.6),
             "សូមអរគុណ", size=28, color=SURFACE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.05), SURFACE)
    add_text(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(0.5),
             "Questions, partnerships, pilots — let's talk.",
             size=18, color=SURFACE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.4),
             "SmartFarm — កសិកម្មឆ្លាតវៃ ងាយស្រួល",
             size=12, color=SURFACE, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_user_flow(prs)
    slide_finance(prs)
    slide_calendar(prs)
    slide_guide(prs)
    slide_settings(prs)
    slide_differentiators(prs)
    slide_status(prs)
    slide_traction(prs)
    slide_competitive(prs)
    slide_team(prs)
    slide_ask(prs)
    slide_appendix_screenshots(prs)
    slide_close(prs)

    out = "SmartFarm-Pitch.pptx"
    prs.save(out)
    print(f"Wrote {out} — {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
